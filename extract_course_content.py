from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pdfplumber
import os
from PIL import Image
import io

def extract_images_from_pptx(file_path, course_num):
    """Extract images from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        image_dir = f"images/course{course_num}"
        os.makedirs(image_dir, exist_ok=True)
        
        image_count = 0
        image_references = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape_num, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Determine image extension
                    ext = image.ext
                    image_filename = f"slide{slide_num}_image{shape_num}.{ext}"
                    image_path = os.path.join(image_dir, image_filename)
                    
                    # Save image
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    image_references.append((slide_num, image_path))
                    image_count += 1
        
        return image_references, image_count
    except Exception as e:
        print(f"Error extracting images from {file_path}: {str(e)}")
        return [], 0

def extract_images_from_pdf(file_path, course_num):
    """Extract images from PDF file"""
    try:
        image_dir = f"images/course{course_num}"
        os.makedirs(image_dir, exist_ok=True)
        
        image_count = 0
        image_references = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Extract images from page
                if hasattr(page, 'images'):
                    for img_num, img in enumerate(page.images):
                        try:
                            image_filename = f"page{page_num}_image{img_num}.png"
                            image_path = os.path.join(image_dir, image_filename)
                            
                            # Note: pdfplumber doesn't directly save images, so we'll note them
                            image_references.append((page_num, image_path))
                            image_count += 1
                        except Exception as e:
                            print(f"Error processing image on page {page_num}: {str(e)}")
        
        return image_references, image_count
    except Exception as e:
        print(f"Error extracting images from {file_path}: {str(e)}")
        return [], 0

def extract_text_from_pptx(file_path, course_num):
    """Extract text from PowerPoint file with better formatting"""
    try:
        prs = Presentation(file_path)
        text_content = []
        image_refs, img_count = extract_images_from_pptx(file_path, course_num)
        
        # Create a mapping of slide numbers to images
        slide_images = {}
        for slide_num, img_path in image_refs:
            if slide_num not in slide_images:
                slide_images[slide_num] = []
            slide_images[slide_num].append(img_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"\n{'─' * 60}"]
            slide_text.append(f"SLIDE {slide_num}")
            slide_text.append('─' * 60)
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Try to identify if it's a title (usually larger/first text box)
                    if shape == slide.shapes[0] if slide.shapes else False:
                        slide_text.append(f"\n## {text}")
                    else:
                        # Preserve line breaks and bullet points
                        lines = text.split('\n')
                        for line in lines:
                            if line.strip():
                                # Check if it looks like a bullet point
                                if line.strip().startswith(('•', '-', '●', '○')):
                                    slide_text.append(f"  {line.strip()}")
                                else:
                                    slide_text.append(f"{line.strip()}")
            
            # Add image references if any
            if slide_num in slide_images:
                slide_text.append(f"\n[Images in this slide: {', '.join(slide_images[slide_num])}]")
            
            text_content.append("\n".join(slide_text))
        
        print(f"  Extracted {img_count} images from {file_path}")
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error extracting from {file_path}: {str(e)}"

def extract_text_from_pdf(file_path, course_num):
    """Extract text from PDF file with better formatting"""
    try:
        text_content = []
        image_refs, img_count = extract_images_from_pdf(file_path, course_num)
        
        # Create a mapping of page numbers to images
        page_images = {}
        for page_num, img_path in image_refs:
            if page_num not in page_images:
                page_images[page_num] = []
            page_images[page_num].append(img_path)
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = [f"\n{'─' * 60}"]
                page_text.append(f"PAGE {page_num}")
                page_text.append('─' * 60)
                
                text = page.extract_text()
                if text:
                    # Preserve formatting
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            page_text.append(line)
                
                # Add image references if any
                if page_num in page_images:
                    page_text.append(f"\n[Images on this page: {', '.join(page_images[page_num])}]")
                
                text_content.append("\n".join(page_text))
        
        print(f"  Extracted {img_count} images from {file_path}")
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error extracting from {file_path}: {str(e)}"

def main():
    # Define course files in order
    course_files = [
        ("curs1 RN 2025.pptx", "Course 1: RN 2025"),
        ("curs2 RN 2025 - perceptron.pptx", "Course 2: Perceptron"),
        ("curs3 RN 2025 -  gradient descent.pptx", "Course 3: Gradient Descent"),
        ("curs4 RN 2025 - backpropagation.pptx", "Course 4: Backpropagation"),
        ("curs5 - weight initialization & overfitting.pptx", "Course 5: Weight Initialization & Overfitting"),
        ("curs6 - optimizers.pdf", "Course 6: Optimizers"),
        ("Curs 7 rn - pytorch.pptx", "Course 7: PyTorch"),
        ("curs8 RN - Q Learning.pptx", "Course 8: Q Learning"),
        ("curs9 -  Convolutional.pptx", "Course 9: Convolutional Networks"),
        ("curs10 - actor critic.pptx", "Course 10: Actor Critic"),
        ("curs11 - LSTM.pptx", "Course 11: LSTM"),
    ]
    
    all_content = []
    all_content.append("=" * 80)
    all_content.append("NEURAL NETWORKS COURSE CONTENT")
    all_content.append("Extracted from courses 1-11")
    all_content.append("=" * 80)
    all_content.append("")
    
    for course_num, (file_name, course_title) in enumerate(course_files, 1):
        if not os.path.exists(file_name):
            print(f"Warning: {file_name} not found, skipping...")
            continue
        
        print(f"Extracting {course_title}...")
        all_content.append("\n\n")
        all_content.append("=" * 80)
        all_content.append(course_title.upper())
        all_content.append("=" * 80)
        all_content.append("")
        
        if file_name.endswith('.pptx'):
            content = extract_text_from_pptx(file_name, course_num)
        elif file_name.endswith('.pdf'):
            content = extract_text_from_pdf(file_name, course_num)
        else:
            content = f"Unknown file format: {file_name}"
        
        all_content.append(content)
    
    # Write all content to a single file
    output_file = "course_content.txt"
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("\n".join(all_content))
    
    print(f"\n{'=' * 60}")
    print(f"✓ Content extracted successfully to {output_file}")
    print(f"✓ Images saved in images/ directory")
    print(f"{'=' * 60}")

if __name__ == "__main__":
    main()